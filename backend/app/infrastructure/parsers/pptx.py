import os
from pptx import Presentation
from app.interfaces.parsers import (
    IDocumentParser, 
    ParsedDocumentResult, 
    ParsedPage, 
    HeadingData, 
    TableData, 
    ImageData, 
    DocumentParsingError
)


class PptxParser(IDocumentParser):
    """Parser implementation for PowerPoint presentations (PPTX)."""

    def parse(self, file_path: str) -> ParsedDocumentResult:
        if not os.path.exists(file_path):
            raise DocumentParsingError(f"File not found: {file_path}")
            
        try:
            prs = Presentation(file_path)
            pages = []
            all_text_parts = []
            
            for idx, slide in enumerate(prs.slides):
                page_number = idx + 1
                slide_text_parts = []
                headings = []
                tables = []
                images = []
                
                # Traverse slide shapes
                for shape_idx, shape in enumerate(slide.shapes):
                    # 1. Text fields
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if not text:
                                continue
                            slide_text_parts.append(text)
                            
                            # Simple slide title heuristic (idx 0 is usually the main slide title)
                            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                                headings.append(HeadingData(text=text, level=1))

                    # 2. Slide tables
                    if shape.has_table:
                        table = shape.table
                        headers = []
                        rows = []
                        for r_idx, row in enumerate(table.rows):
                            row_cells = [cell.text.strip() for cell in row.cells]
                            if r_idx == 0:
                                headers = row_cells
                            else:
                                rows.append(row_cells)
                                
                        tables.append(TableData(
                            headers=headers,
                            rows=rows,
                            caption=f"Slide {page_number} Table {shape_idx + 1}"
                        ))
                        # Append table rows to search index text representation
                        for row_data in rows:
                            slide_text_parts.append(" | ".join(row_data))

                    # 3. Slide embedded pictures
                    if hasattr(shape, "image") and shape.image:
                        try:
                            image = shape.image
                            ext = image.ext or "png"
                            images.append(ImageData(
                                name=f"slide{page_number}_img{shape_idx}.{ext}",
                                content=image.blob,
                                content_type=f"image/{ext}",
                                page_number=page_number
                            ))
                        except Exception:
                            pass

                slide_text = "\n".join(slide_text_parts)
                pages.append(ParsedPage(
                    page_number=page_number,
                    text=slide_text,
                    headings=headings,
                    tables=tables,
                    images=images
                ))
                if slide_text:
                    all_text_parts.append(slide_text)

            raw_text = "\n\n".join(all_text_parts)
            
            # Extract pptx core metadata properties
            metadata = {}
            try:
                props = prs.core_properties
                for prop in ["author", "created", "last_modified_by", "modified", "title"]:
                    val = getattr(props, prop)
                    if val:
                        metadata[prop] = str(val)
            except Exception:
                pass
            metadata["page_count"] = len(prs.slides)

            return ParsedDocumentResult(
                raw_text=raw_text,
                metadata=metadata,
                pages=pages
            )
        except Exception as e:
            raise DocumentParsingError(f"Failed parsing PPTX file: {str(e)}")
