import os
import tempfile
import pytest
import csv
import json
from docx import Document as DocxDoc
from openpyxl import Workbook
from pptx import Presentation

from app.interfaces.parsers import DocumentParsingError, UnsupportedFormatError
from app.infrastructure.parsers.factory import DocumentParserFactory
from app.infrastructure.parsers.pdf import PDFParser
from app.infrastructure.parsers.docx import DocxParser
from app.infrastructure.parsers.excel import ExcelParser
from app.infrastructure.parsers.pptx import PptxParser
from app.infrastructure.parsers.markdown import MarkdownParser
from app.infrastructure.parsers.html import HtmlParser
from app.infrastructure.parsers.csv import CsvParser
from app.infrastructure.parsers.json import JsonParser
from app.infrastructure.parsers.txt import TxtParser


# -------------------------------------------------------------
# Test Parser Factory
# -------------------------------------------------------------

def test_parser_factory_resolution():
    """Verify that DocumentParserFactory resolves correct instances or raises errors."""
    factory = DocumentParserFactory()
    
    assert isinstance(factory.get_parser("pdf"), PDFParser)
    assert isinstance(factory.get_parser(".docx"), DocxParser)
    assert isinstance(factory.get_parser("XLSX"), ExcelParser)
    assert isinstance(factory.get_parser("ppt"), PptxParser)
    
    with pytest.raises(UnsupportedFormatError):
        factory.get_parser("png")


# -------------------------------------------------------------
# Test Text, CSV, JSON, HTML, Markdown Parsers
# -------------------------------------------------------------

def test_txt_parser():
    """Verify plain text parser counts characters, words, and lines correctly."""
    parser = TxtParser()
    content = "Hello line 1\nSecond line with some words."
    
    with tempfile.NamedTemporaryFile(mode="w", delete=False, encoding="utf-8") as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        assert res.raw_text == content
        assert res.metadata["line_count"] == "2"
        assert res.metadata["word_count"] == "8"
        assert len(res.pages) == 1
    finally:
        os.remove(tmp_path)


def test_txt_parser_corrupted():
    """Verify TxtParser throws DocumentParsingError on invalid UTF-8 files."""
    parser = TxtParser()
    with tempfile.NamedTemporaryFile(mode="wb", delete=False) as tmp:
        # Write invalid UTF-8 bytes (binary payload)
        tmp.write(b"\xff\xfe\x00\x00\x01\x02")
        tmp_path = tmp.name

    try:
        with pytest.raises(DocumentParsingError) as excinfo:
            parser.parse(tmp_path)
        assert "decoding failed" in str(excinfo.value)
    finally:
        os.remove(tmp_path)


def test_csv_parser():
    """Verify CSV parser creates structured table outputs."""
    parser = CsvParser()
    
    with tempfile.NamedTemporaryFile(mode="w", newline="", delete=False, encoding="utf-8") as tmp:
        writer = csv.writer(tmp)
        writer.writerow(["Name", "Age", "Role"])
        writer.writerow(["Alice", "30", "Engineer"])
        writer.writerow(["Bob", "25", "Designer"])
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        assert len(res.pages) == 1
        page = res.pages[0]
        assert len(page.tables) == 1
        table = page.tables[0]
        assert table.headers == ["Name", "Age", "Role"]
        assert table.rows == [["Alice", "30", "Engineer"], ["Bob", "25", "Designer"]]
        assert res.metadata["row_count"] == "3"
        assert "Alice | 30 | Engineer" in res.raw_text
    finally:
        os.remove(tmp_path)


def test_json_parser():
    """Verify JSON parser extracts formatted text dumps and structure details."""
    parser = JsonParser()
    data = {"name": "Test Org", "users": ["User A", "User B"]}
    
    with tempfile.NamedTemporaryFile(mode="w", delete=False, encoding="utf-8") as tmp:
        json.dump(data, tmp)
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        parsed_data = json.loads(res.raw_text)
        assert parsed_data == data
        assert res.metadata["json_type"] == "object"
        assert res.metadata["item_count"] == "2"
    finally:
        os.remove(tmp_path)


def test_html_parser():
    """Verify HTML parser extracts tags, text, headings, and tables."""
    parser = HtmlParser()
    html_content = """
    <html>
        <head><title>My Test Page</title></head>
        <body>
            <h1>Heading 1 text</h1>
            <p>Welcome to parsing!</p>
            <table>
                <caption>Staff list</caption>
                <tr><th>User</th><th>Active</th></tr>
                <tr><td>Alice</td><td>Yes</td></tr>
            </table>
        </body>
    </html>
    """
    
    with tempfile.NamedTemporaryFile(mode="w", delete=False, encoding="utf-8") as tmp:
        tmp.write(html_content)
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        assert res.metadata["title"] == "My Test Page"
        assert len(res.pages) == 1
        page = res.pages[0]
        
        # Headings
        assert len(page.headings) == 1
        assert page.headings[0].text == "Heading 1 text"
        assert page.headings[0].level == 1
        
        # Tables
        assert len(page.tables) == 1
        assert page.tables[0].headers == ["User", "Active"]
        assert page.tables[0].rows == [["Alice", "Yes"]]
        assert page.tables[0].caption == "Staff list"
    finally:
        os.remove(tmp_path)


def test_markdown_parser():
    """Verify markdown parses headers, text, and tables."""
    parser = MarkdownParser()
    md_content = """# Title Header
Some text.

| Name | Role |
| --- | --- |
| Alice | Admin |
"""
    
    with tempfile.NamedTemporaryFile(mode="w", delete=False, encoding="utf-8") as tmp:
        tmp.write(md_content)
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        assert res.metadata["title"] == "Title Header"
        page = res.pages[0]
        assert len(page.headings) == 1
        assert page.headings[0].text == "Title Header"
        
        assert len(page.tables) == 1
        assert page.tables[0].headers == ["Name", "Role"]
        assert page.tables[0].rows == [["Alice", "Admin"]]
    finally:
        os.remove(tmp_path)


# -------------------------------------------------------------
# Test Heavy Binary Parsers (Word, Excel, PPTX, PDF)
# -------------------------------------------------------------

def test_docx_parser():
    """Verify DocxParser parses structured paragraphs and tables."""
    parser = DocxParser()
    
    doc = DocxDoc()
    doc.add_heading("Docx Title", level=2)
    doc.add_paragraph("General body content.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Col A"
    table.cell(0, 1).text = "Col B"
    table.cell(1, 0).text = "Val A"
    table.cell(1, 1).text = "Val B"
    
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        doc.save(tmp.name)
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        page = res.pages[0]
        assert len(page.headings) == 1
        assert page.headings[0].text == "Docx Title"
        assert page.headings[0].level == 2
        
        assert len(page.tables) == 1
        assert page.tables[0].headers == ["Col A", "Col B"]
        assert page.tables[0].rows == [["Val A", "Val B"]]
    finally:
        os.remove(tmp_path)


def test_docx_parser_corrupted():
    """Verify DocxParser throws DocumentParsingError on malformed documents."""
    parser = DocxParser()
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp.write(b"NOT A ZIP OR XML DOCX STRUCTURE")
        tmp_path = tmp.name

    try:
        with pytest.raises(DocumentParsingError):
            parser.parse(tmp_path)
    finally:
        os.remove(tmp_path)


def test_excel_parser():
    """Verify ExcelParser translates spreadsheets to tables and logical pages."""
    parser = ExcelParser()
    
    wb = Workbook()
    ws = wb.active
    ws.title = "SummarySheet"
    ws.append(["Item", "Price"])
    ws.append(["Laptop", 1200])
    
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
        wb.save(tmp.name)
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        assert res.metadata["sheet_count"] == 1
        page = res.pages[0]
        assert len(page.tables) == 1
        assert page.tables[0].headers == ["Item", "Price"]
        assert page.tables[0].rows == [["Laptop", "1200"]]
    finally:
        os.remove(tmp_path)


def test_pptx_parser():
    """Verify PptxParser slide traversal."""
    parser = PptxParser()
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Presentation Title"
    
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp_path = tmp.name

    try:
        res = parser.parse(tmp_path)
        assert len(res.pages) == 1
        page = res.pages[0]
        assert page.headings[0].text == "Presentation Title"
        assert "Presentation Title" in page.text
    finally:
        os.remove(tmp_path)


def test_pdf_parser_corrupted():
    """Verify PDFParser raises DocumentParsingError on empty/corrupted files."""
    parser = PDFParser()
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(b"CORRUPT DATA STREAM")
        tmp_path = tmp.name

    try:
        with pytest.raises(DocumentParsingError) as excinfo:
            parser.parse(tmp_path)
        # Verify it mapped the exception
        assert "Failed parsing PDF file" in str(excinfo.value)
    finally:
        os.remove(tmp_path)
